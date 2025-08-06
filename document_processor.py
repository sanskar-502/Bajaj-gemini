import os
import uuid
from datetime import datetime
from typing import List, Dict, Any, Optional
import PyPDF2
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("Warning: PyMuPDF not available, using PyPDF2 fallback for PDF processing")
from docx import Document
import unstructured
from unstructured.partition.auto import partition
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
import re

from config import Config
from models import DocumentMetadata, DocumentType

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

class DocumentProcessor:
    """Handles document parsing, chunking, and metadata extraction"""
    
    def __init__(self):
        self.supported_formats = Config.SUPPORTED_FORMATS
        self.chunk_size = Config.CHUNK_SIZE
        self.chunk_overlap = Config.CHUNK_OVERLAP
    
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Process a document and return chunks with metadata
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing chunks, metadata, and processing info
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        # Extract text and metadata
        text_content, metadata = self._extract_text_and_metadata(file_path)
        print(f"[DocumentProcessor] Extracted text length: {len(text_content)}")
        if not text_content.strip():
            print(f"[DocumentProcessor] Warning: No text extracted from {file_path}")
        
        # Create chunks
        chunks = self._create_chunks(text_content, metadata)
        print(f"[DocumentProcessor] Chunks created: {len(chunks)}")
        
        # Generate document ID
        document_id = str(uuid.uuid4())
        
        return {
            "document_id": document_id,
            "chunks": chunks,
            "metadata": metadata,
            "total_chunks": len(chunks)
        }
    
    def _extract_text_and_metadata(self, file_path: str) -> tuple[str, DocumentMetadata]:
        """Extract text content and metadata from document"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == ".pdf":
            return self._process_pdf(file_path)
        elif file_extension == ".docx":
            return self._process_docx(file_path)
        elif file_extension == ".txt":
            return self._process_txt(file_path)
        elif file_extension == ".pptx":
            return self._process_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _process_pdf(self, file_path: str) -> tuple[str, DocumentMetadata]:
        """Process PDF document"""
        text_content = ""
        page_count = 0
        
        # Try PyMuPDF first (better text extraction) if available
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                page_count = len(doc)
                
                for page_num in range(page_count):
                    page = doc.load_page(page_num)
                    text_content += page.get_text()
                
                doc.close()
            except Exception as e:
                print(f"PyMuPDF failed, falling back to PyPDF2: {str(e)}")
                # Fallback to PyPDF2
                try:
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        page_count = len(pdf_reader.pages)
                        
                        for page in pdf_reader.pages:
                            text_content += page.extract_text()
                except Exception as e2:
                    raise Exception(f"Failed to process PDF: {str(e2)}")
        else:
            # Use PyPDF2 directly
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    page_count = len(pdf_reader.pages)
                    
                    for page in pdf_reader.pages:
                        text_content += page.extract_text()
            except Exception as e:
                raise Exception(f"Failed to process PDF: {str(e)}")
        
        # Clean text
        text_content = self._clean_text(text_content)
        
        # Detect document type
        doc_type = self._detect_document_type(text_content)
        
        # Extract company name
        company_name = self._extract_company_name(text_content)
        
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            document_type=doc_type,
            upload_timestamp=datetime.now().isoformat(),
            file_size=os.path.getsize(file_path),
            page_count=page_count,
            company_name=company_name,
            document_version=None
        )
        
        return text_content, metadata
    
    def _process_docx(self, file_path: str) -> tuple[str, DocumentMetadata]:
        """Process DOCX document"""
        doc = Document(file_path)
        text_content = ""
        for paragraph in doc.paragraphs:
            text_content += paragraph.text + "\n"
        print(f"[DocumentProcessor] DOCX extracted text (first 300 chars): {text_content[:300]}")
        # Clean text
        text_content = self._clean_text(text_content)
        if not text_content.strip():
            print(f"[DocumentProcessor] Warning: Cleaned DOCX text is empty for {file_path}")
        # Detect document type
        doc_type = self._detect_document_type(text_content)
        # Extract company name
        company_name = self._extract_company_name(text_content)
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            document_type=doc_type,
            upload_timestamp=datetime.now().isoformat(),
            file_size=os.path.getsize(file_path),
            page_count=None,
            company_name=company_name,
            document_version=None
        )
        return text_content, metadata
    
    def _process_txt(self, file_path: str) -> tuple[str, DocumentMetadata]:
        """Process TXT document"""
        with open(file_path, 'r', encoding='utf-8') as file:
            text_content = file.read()
        
        # Clean text
        text_content = self._clean_text(text_content)
        
        # Detect document type
        doc_type = self._detect_document_type(text_content)
        
        # Extract company name
        company_name = self._extract_company_name(text_content)
        
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            document_type=doc_type,
            upload_timestamp=datetime.now().isoformat(),
            file_size=os.path.getsize(file_path),
            page_count=None,
            company_name=company_name,
            document_version=None
        )
        
        return text_content, metadata
    
    def _process_pptx(self, file_path: str) -> tuple[str, DocumentMetadata]:
        """Process PPTX document"""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_content = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
        
        # Clean text
        text_content = self._clean_text(text_content)
        
        # Detect document type
        doc_type = self._detect_document_type(text_content)
        
        # Extract company name
        company_name = self._extract_company_name(text_content)
        
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            document_type=doc_type,
            upload_timestamp=datetime.now().isoformat(),
            file_size=os.path.getsize(file_path),
            page_count=len(prs.slides),
            company_name=company_name,
            document_version=None
        )
        
        return text_content, metadata
    
    def _create_chunks(self, text_content: str, metadata: DocumentMetadata) -> List[Dict[str, Any]]:
        """Create chunks from text content"""
        chunks = []
        
        # Split into sentences first
        sentences = sent_tokenize(text_content)
        if not sentences:
            print("[DocumentProcessor] Warning: sent_tokenize returned no sentences, falling back to line-based chunking.")
            sentences = text_content.splitlines()
        
        current_chunk = ""
        chunk_id = 0
        
        for sentence in sentences:
            # If adding this sentence would exceed chunk size, save current chunk
            if len(current_chunk) + len(sentence) > self.chunk_size and current_chunk:
                chunks.append(self._create_chunk_dict(current_chunk, chunk_id, metadata))
                chunk_id += 1
                
                # Start new chunk with overlap
                overlap_sentences = current_chunk.split('.')[-3:]  # Last 3 sentences
                current_chunk = '. '.join(overlap_sentences) + '. ' + sentence
            else:
                current_chunk += sentence + " "
        
        # Add the last chunk if it has content
        if current_chunk.strip():
            chunks.append(self._create_chunk_dict(current_chunk, chunk_id, metadata))
        
        return chunks
    
    def _create_chunk_dict(self, text: str, chunk_id: int, metadata: DocumentMetadata) -> Dict[str, Any]:
        """Create a chunk dictionary with metadata"""
        return {
            "id": f"{metadata.filename}_{chunk_id}",
            "text": text.strip(),
            "metadata": {
                "document_id": metadata.filename,
                "chunk_id": chunk_id,
                "document_type": metadata.document_type.value,
                "company_name": metadata.company_name,
                "page_count": metadata.page_count,
                "upload_timestamp": metadata.upload_timestamp
            }
        }
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content"""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)\[\]\{\}]', ' ', text)
        
        # Normalize quotes
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace(''', "'").replace(''', "'")
        
        return text.strip()
    
    def _detect_document_type(self, text: str) -> DocumentType:
        """Detect document type based on content"""
        text_lower = text.lower()
        
        # Insurance policy keywords
        insurance_keywords = ['policy', 'coverage', 'premium', 'claim', 'insured', 'insurance', 'deductible']
        if any(keyword in text_lower for keyword in insurance_keywords):
            return DocumentType.INSURANCE_POLICY
        
        # Legal contract keywords
        legal_keywords = ['contract', 'agreement', 'terms', 'conditions', 'party', 'obligation', 'liability']
        if any(keyword in text_lower for keyword in legal_keywords):
            return DocumentType.LEGAL_CONTRACT
        
        # HR policy keywords
        hr_keywords = ['employee', 'employment', 'termination', 'leave', 'benefits', 'conduct', 'harassment']
        if any(keyword in text_lower for keyword in hr_keywords):
            return DocumentType.HR_POLICY
        
        # Compliance keywords
        compliance_keywords = ['compliance', 'regulation', 'gdpr', 'hipaa', 'sox', 'audit', 'policy']
        if any(keyword in text_lower for keyword in compliance_keywords):
            return DocumentType.COMPLIANCE_DOC
        
        return DocumentType.UNKNOWN
    
    def _extract_company_name(self, text: str) -> Optional[str]:
        """Extract company name from document content"""
        # Look for common patterns
        patterns = [
            r'(?:company|corporation|inc|llc|ltd)\s*[:=]\s*([A-Z][A-Za-z\s&]+)',
            r'([A-Z][A-Za-z\s&]+)\s*(?:company|corporation|inc|llc|ltd)',
            r'between\s+([A-Z][A-Za-z\s&]+)\s+and',
            r'this\s+agreement\s+is\s+between\s+([A-Z][A-Za-z\s&]+)'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                company_name = match.group(1).strip()
                if len(company_name) > 2:  # Filter out very short matches
                    return company_name
        
        return None 